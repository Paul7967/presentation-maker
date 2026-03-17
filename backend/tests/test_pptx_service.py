"""Unit tests for pptx_service: parsing template and generating presentation."""

import io

import pytest
from pptx import Presentation

from app.services.pptx_service import (
    MAX_ITEMS_PER_SLIDE,
    MAX_SLIDES,
    extract_style_from_template,
    generate_presentation,
)


class TestExtractStyleFromTemplate:
    """Tests for extract_style_from_template."""

    def test_extract_style_returns_layout_and_styles(self, minimal_pptx_bytes):
        """Parse template and get layout, title_style, body_style from first slide."""
        layout, title_style, body_style = extract_style_from_template(minimal_pptx_bytes)
        assert layout is not None
        # Styles may be None if placeholders are empty; at least no exception
        assert title_style is None or isinstance(title_style, dict)
        assert body_style is None or isinstance(body_style, dict)

    def test_extract_style_empty_pptx_raises(self):
        """Template with no slides raises ValueError('no_slides')."""
        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        empty_bytes = buf.read()
        with pytest.raises(ValueError, match="no_slides"):
            extract_style_from_template(empty_bytes)

    def test_extract_style_invalid_bytes_raises(self):
        """Invalid bytes raise (e.g. not a valid pptx)."""
        with pytest.raises(Exception):  # python-pptx may raise various exceptions
            extract_style_from_template(b"not a zip file")


class TestGeneratePresentation:
    """Tests for generate_presentation."""

    def test_generate_returns_bytes(self, minimal_pptx_bytes):
        """Generate returns bytes (valid pptx)."""
        slides_data = [
            {"title": "Slide 1", "items": ["Point A", "Point B"]},
            {"title": "Slide 2", "items": []},
        ]
        result = generate_presentation(minimal_pptx_bytes, slides_data)
        assert isinstance(result, bytes)
        assert len(result) > 0
        assert result[:2] == b"PK"

    def test_generate_slide_count_matches_input(self, minimal_pptx_bytes):
        """Output pptx has the same number of slides as slides_data."""
        slides_data = [
            {"title": f"Slide {i}", "items": [f"Item {i}"]}
            for i in range(1, 6)
        ]
        result = generate_presentation(minimal_pptx_bytes, slides_data)
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 5

    def test_generate_one_slide(self, minimal_pptx_bytes):
        """Single slide: overwrites first slide of template."""
        slides_data = [{"title": "Only", "items": ["One item"]}]
        result = generate_presentation(minimal_pptx_bytes, slides_data)
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 1

    def test_generate_empty_slides_raises(self, minimal_pptx_bytes):
        """Empty slides_data raises ValueError."""
        with pytest.raises(ValueError, match="invalid_slide_count"):
            generate_presentation(minimal_pptx_bytes, [])

    def test_generate_too_many_slides_raises(self, minimal_pptx_bytes):
        """More than MAX_SLIDES raises ValueError."""
        slides_data = [
            {"title": "", "items": []}
            for _ in range(MAX_SLIDES + 1)
        ]
        with pytest.raises(ValueError, match="invalid_slide_count"):
            generate_presentation(minimal_pptx_bytes, slides_data)

    def test_generate_empty_title_and_items_allowed(self, minimal_pptx_bytes):
        """Slide with empty title and empty items is valid."""
        slides_data = [{"title": "", "items": []}]
        result = generate_presentation(minimal_pptx_bytes, slides_data)
        assert isinstance(result, bytes)
        assert len(result) > 0

    def test_generate_with_layout_id_index(self, minimal_pptx_bytes):
        """Slides with layoutId (index) use that layout for added slides; fallback for first."""
        slides_data = [
            {"title": "First", "items": [], "layoutId": 0},
            {"title": "Second", "items": [], "layoutId": 1},
        ]
        result = generate_presentation(minimal_pptx_bytes, slides_data)
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 2

    def test_generate_with_notes(self, minimal_pptx_bytes):
        """Slides with notes field write to notes_slide."""
        slides_data = [
            {"title": "Slide 1", "items": [], "notes": "First speaker note"},
            {"title": "Slide 2", "items": [], "notes": ""},
        ]
        result = generate_presentation(minimal_pptx_bytes, slides_data)
        assert isinstance(result, bytes)
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 2
        # First slide has notes_slide; notes_text_frame may exist
        notes_slide = prs.slides[0].notes_slide
        assert notes_slide is not None

    def test_generate_many_items_per_slide_truncated(self, minimal_pptx_bytes):
        """More than MAX_ITEMS_PER_SLIDE: service truncates (sliced in code)."""
        many_items = [f"Item {i}" for i in range(MAX_ITEMS_PER_SLIDE + 5)]
        slides_data = [{"title": "Many", "items": many_items}]
        result = generate_presentation(minimal_pptx_bytes, slides_data)
        assert isinstance(result, bytes)
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 1
