"""Integration tests for POST /api/presentation/generate."""

import io
import json
import os
import tempfile

import pytest
from fastapi.testclient import TestClient

from app.main import app

client = TestClient(app)

# Minimal valid pptx (PK + zip structure with ppt/)
MINIMAL_PPTX = (
    b"PK\x03\x04"
    b"\x14\x00\x00\x00\x08\x00\x00\x00!\x00\x00\x00[Content_Types].xml \x00\x00\x00"
)
# Ensure ppt/ appears in first 4096 bytes for _is_pptx_signature
MINIMAL_PPTX_BODY = MINIMAL_PPTX + b"ppt/" + b"\x00" * 100


def _make_pptx_file():
    """Create a minimal valid pptx file content for upload (real pptx via python-pptx in conftest)."""
    from pptx import Presentation
    buf = io.BytesIO()
    prs = Presentation()
    layout = prs.slide_layouts[0]
    prs.slides.add_slide(layout)
    prs.save(buf)
    buf.seek(0)
    return buf.read()


@pytest.fixture
def sample_pptx_bytes():
    """Valid pptx bytes for multipart upload."""
    return _make_pptx_file()


@pytest.fixture
def valid_slides_json():
    """Valid slides JSON string."""
    return json.dumps([
        {"title": "First", "items": ["A", "B"]},
        {"title": "Second", "items": []},
    ])


class TestGenerateSuccess:
    """Successful generation returns 200 and binary pptx."""

    def test_200_with_pptx(self, sample_pptx_bytes, valid_slides_json):
        """Upload template + slides -> 200, Content-Type pptx, binary body."""
        files = {"template": ("template.pptx", sample_pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
        data = {"slides": valid_slides_json}
        response = client.post("/api/presentation/generate", data=data, files=files)
        assert response.status_code == 200
        assert "application/vnd.openxmlformats-officedocument.presentationml.presentation" in (
            response.headers.get("content-type", "")
        )
        assert "attachment" in response.headers.get("content-disposition", "").lower()
        body = response.content
        assert len(body) > 0
        assert body[:2] == b"PK"

    def test_response_size_reasonable(self, sample_pptx_bytes, valid_slides_json):
        """Response body is non-empty and looks like pptx."""
        files = {"template": ("t.pptx", sample_pptx_bytes, "application/octet-stream")}
        data = {"slides": valid_slides_json}
        response = client.post("/api/presentation/generate", data=data, files=files)
        assert response.status_code == 200
        assert len(response.content) > 100


class TestGenerateErrors400:
    """Client errors: 400 Bad Request."""

    def test_no_template(self, valid_slides_json):
        """Missing template -> 400 or 422 (FastAPI returns 422 for missing File)."""
        data = {"slides": valid_slides_json}
        response = client.post("/api/presentation/generate", data=data)
        assert response.status_code in (400, 422)
        detail = response.json().get("detail", "")
        if isinstance(detail, list):
            detail = " ".join(str(d.get("msg", d)) for d in detail)
        assert "образец" in detail.lower() or "pptx" in detail.lower() or "template" in detail.lower() or "required" in detail.lower()

    def test_no_slides(self, sample_pptx_bytes):
        """Missing slides field -> 400 or 422 (FastAPI returns 422 for missing Form)."""
        files = {"template": ("t.pptx", sample_pptx_bytes, "application/octet-stream")}
        response = client.post("/api/presentation/generate", files=files)
        assert response.status_code in (400, 422)

    def test_empty_slides_string(self, sample_pptx_bytes):
        """Empty string for slides -> 400 or 422."""
        files = {"template": ("t.pptx", sample_pptx_bytes, "application/octet-stream")}
        data = {"slides": ""}
        response = client.post("/api/presentation/generate", data=data, files=files)
        assert response.status_code in (400, 422)

    def test_slides_not_json(self, sample_pptx_bytes):
        """Invalid JSON in slides -> 400."""
        files = {"template": ("t.pptx", sample_pptx_bytes, "application/octet-stream")}
        data = {"slides": "not json"}
        response = client.post("/api/presentation/generate", data=data, files=files)
        assert response.status_code == 400
        assert "Некорректный" in response.json().get("detail", "") or "формат" in response.json().get("detail", "").lower()

    def test_slides_empty_array(self, sample_pptx_bytes):
        """slides = [] -> 400."""
        files = {"template": ("t.pptx", sample_pptx_bytes, "application/octet-stream")}
        data = {"slides": "[]"}
        response = client.post("/api/presentation/generate", data=data, files=files)
        assert response.status_code == 400
        assert "хотя бы одного слайда" in response.json().get("detail", "")

    def test_slides_more_than_100(self, sample_pptx_bytes):
        """More than 100 slides -> 400."""
        slides = [{"title": "", "items": []} for _ in range(101)]
        files = {"template": ("t.pptx", sample_pptx_bytes, "application/octet-stream")}
        data = {"slides": json.dumps(slides)}
        response = client.post("/api/presentation/generate", data=data, files=files)
        assert response.status_code == 400
        assert "100" in response.json().get("detail", "")

    def test_file_not_pptx_signature(self, valid_slides_json):
        """File that does not have ppt/ signature -> 400."""
        files = {"template": ("x.txt", b"plain text file", "application/octet-stream")}
        data = {"slides": valid_slides_json}
        response = client.post("/api/presentation/generate", data=data, files=files)
        assert response.status_code == 400
        assert "PPTX" in response.json().get("detail", "") or "формат" in response.json().get("detail", "").lower()

    def test_empty_file(self, valid_slides_json):
        """Empty template file -> 400."""
        files = {"template": ("e.pptx", b"", "application/octet-stream")}
        data = {"slides": valid_slides_json}
        response = client.post("/api/presentation/generate", data=data, files=files)
        assert response.status_code == 400


class TestGenerateError413:
    """Payload too large: 413."""

    def test_content_length_over_20mb_returns_413(self, sample_pptx_bytes, valid_slides_json):
        """Request with body size > 20 MB returns 413 (middleware or endpoint)."""
        size_over = 20 * 1024 * 1024 + 1
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
                tmp_path = f.name
                f.write(sample_pptx_bytes)
                f.write(b"\x00" * (size_over - len(sample_pptx_bytes)))
                f.flush()
            with open(tmp_path, "rb") as r:
                big_body = r.read()
            files = {"template": ("big.pptx", big_body, "application/octet-stream")}
            data = {"slides": valid_slides_json}
            response = client.post("/api/presentation/generate", data=data, files=files)
            assert response.status_code == 413
            assert "20" in response.json().get("detail", "") and "MB" in response.json().get("detail", "")
        finally:
            if tmp_path and os.path.isfile(tmp_path):
                os.unlink(tmp_path)


class TestGenerateValidationLimits:
    """Validation limits: title length, item length, items per slide."""

    def test_title_too_long_400(self, sample_pptx_bytes):
        """Title > 500 chars -> 400."""
        long_title = "x" * 501
        slides = [{"title": long_title, "items": []}]
        files = {"template": ("t.pptx", sample_pptx_bytes, "application/octet-stream")}
        data = {"slides": json.dumps(slides)}
        response = client.post("/api/presentation/generate", data=data, files=files)
        assert response.status_code == 400
        assert "500" in response.json().get("detail", "")

    def test_item_too_long_400(self, sample_pptx_bytes):
        """Item > 2000 chars -> 400."""
        long_item = "y" * 2001
        slides = [{"title": "", "items": [long_item]}]
        files = {"template": ("t.pptx", sample_pptx_bytes, "application/octet-stream")}
        data = {"slides": json.dumps(slides)}
        response = client.post("/api/presentation/generate", data=data, files=files)
        assert response.status_code == 400
        assert "2000" in response.json().get("detail", "") or "пункт" in response.json().get("detail", "")

    def test_too_many_items_per_slide_400(self, sample_pptx_bytes):
        """More than 50 items on one slide -> 400."""
        items = ["z"] * 51
        slides = [{"title": "", "items": items}]
        files = {"template": ("t.pptx", sample_pptx_bytes, "application/octet-stream")}
        data = {"slides": json.dumps(slides)}
        response = client.post("/api/presentation/generate", data=data, files=files)
        assert response.status_code == 400
        assert "50" in response.json().get("detail", "")
