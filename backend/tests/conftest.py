"""Pytest configuration and shared fixtures."""

import io

import pytest
from pptx import Presentation


@pytest.fixture
def minimal_pptx_bytes():
    """
    Create a minimal valid PPTX (one slide, Title and Content layout) in memory.
    Used for tests that need a real template.
    """
    prs = Presentation()
    # Use built-in layout index 0 (Title and Content)
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    title_ph = slide.shapes.title
    if title_ph:
        title_ph.text = "Sample Title"
    # Body placeholder is typically index 1
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body = shape
            break
    if body and body.has_text_frame:
        body.text_frame.text = "Sample body"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
