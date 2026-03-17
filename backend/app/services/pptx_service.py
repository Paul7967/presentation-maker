"""Service for parsing template PPTX and generating new presentations."""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation


# Placeholder indices: 0 = title, 1 = body (content) in standard Title and Content layout
TITLE_PLACEHOLDER_IDX = 0
BODY_PLACEHOLDER_IDX = 1

# Limits from api.md
MAX_TITLE_LEN = 500
MAX_ITEM_LEN = 2000
MAX_ITEMS_PER_SLIDE = 50
MAX_SLIDES = 100
MAX_NOTES_LEN = 4000
MAX_FILE_SIZE = 20 * 1024 * 1024  # 20 MB


def _get_placeholder(slide, idx: int):
    """Get placeholder by index if present."""
    try:
        return slide.placeholders[idx]
    except KeyError:
        return None


def _font_info(shape) -> dict[str, Any] | None:
    """Extract font info from first paragraph of a text frame."""
    if not shape or not getattr(shape, "has_text_frame", True):
        return None
    if not shape.has_text_frame:
        return None
    tf = shape.text_frame
    if not tf.paragraphs:
        return None
    p = tf.paragraphs[0]
    if not p.runs:
        return None
    r = p.runs[0]
    font = r.font
    info = {
        "name": font.name,
        "size": font.size,
        "bold": font.bold,
        "italic": font.italic,
    }
    if getattr(font.color, "rgb", None) is not None:
        info["rgb"] = font.color.rgb
    return info


def extract_style_from_template(template_bytes: bytes) -> tuple[Any, dict | None, dict | None]:
    """
    Parse template PPTX and extract style from the first slide.
    Returns (layout, title_style_dict, body_style_dict).
    Raises ValueError("no_slides") if template has no slides.
    """
    stream = io.BytesIO(template_bytes)
    prs = Presentation(stream)
    if not prs.slides:
        raise ValueError("no_slides")
    slide = prs.slides[0]
    layout = slide.slide_layout
    title_placeholder = _get_placeholder(slide, TITLE_PLACEHOLDER_IDX)
    body_placeholder = _get_placeholder(slide, BODY_PLACEHOLDER_IDX)
    title_style = _font_info(title_placeholder) if title_placeholder else None
    body_style = _font_info(body_placeholder) if body_placeholder else None
    return (layout, title_style, body_style)


def _set_placeholder_text(placeholder, lines: list[str], font_style: dict | None):
    """Set text in a placeholder (title = single line, body = list of items)."""
    if not placeholder or not placeholder.has_text_frame:
        return
    tf = placeholder.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        else:
            p = tf.add_paragraph()
        p.text = (line or "")[: (MAX_ITEM_LEN if len(lines) > 1 else MAX_TITLE_LEN)]
        p.level = 0
        if font_style and p.runs:
            r = p.runs[0]
            if font_style.get("name"):
                r.font.name = font_style["name"]
            if font_style.get("size") is not None:
                r.font.size = font_style["size"]
            if font_style.get("bold") is not None:
                r.font.bold = font_style["bold"]
            if font_style.get("italic") is not None:
                r.font.italic = font_style["italic"]
            if font_style.get("rgb") is not None:
                r.font.color.rgb = font_style["rgb"]


def _set_notes_text(slide, notes_text: str) -> None:
    """Write notes text to slide's notes_slide. No-op if notes_text_frame is missing."""
    if not notes_text or not slide:
        return
    text = (notes_text or "").strip()[:MAX_NOTES_LEN]
    if not text:
        return
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        if tf is None:
            return
        tf.clear()
        if tf.paragraphs:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
    except Exception:
        pass


def generate_presentation(template_bytes: bytes, slides_data: list[dict]) -> bytes:
    """
    Generate a new PPTX from template and list of slides.
    slides_data: list of {"title": str, "items": list[str], "layoutId": int|str|None, "notes": str|None}
    Uses template's first slide for slide 0; adds remaining slides with layout per layoutId (fallback: first layout).
    """
    if not slides_data or len(slides_data) > MAX_SLIDES:
        raise ValueError("invalid_slide_count")
    stream = io.BytesIO(template_bytes)
    prs = Presentation(stream)
    if not prs.slides:
        raise ValueError("no_slides")
    layouts = prs.slide_layouts
    first_slide = prs.slides[0]
    title_ph = _get_placeholder(first_slide, TITLE_PLACEHOLDER_IDX)
    body_ph = _get_placeholder(first_slide, BODY_PLACEHOLDER_IDX)
    title_style = _font_info(title_ph) if title_ph else None
    body_style = _font_info(body_ph) if body_ph else None

    def get_layout(layout_id: int | str | None):
        if layout_id is None:
            return layouts[0]
        if isinstance(layout_id, int):
            if 0 <= layout_id < len(layouts):
                return layouts[layout_id]
            return layouts[0]
        if isinstance(layout_id, str):
            layout = layouts.get_by_name(layout_id, None)
            return layout if layout is not None else layouts[0]
        return layouts[0]

    def fill_slide(slide, slide_spec: dict):
        title_ph = _get_placeholder(slide, TITLE_PLACEHOLDER_IDX)
        body_ph = _get_placeholder(slide, BODY_PLACEHOLDER_IDX)
        title_text = (slide_spec.get("title") or "")[:MAX_TITLE_LEN]
        items = (slide_spec.get("items") or [])[:MAX_ITEMS_PER_SLIDE]
        if title_ph:
            _set_placeholder_text(title_ph, [title_text], title_style)
        if body_ph:
            _set_placeholder_text(body_ph, items, body_style)
        _set_notes_text(slide, slide_spec.get("notes"))

    # First slide: overwrite content and notes (keep template layout)
    fill_slide(prs.slides[0], slides_data[0])

    # Add remaining slides with chosen layout
    for slide_spec in slides_data[1:]:
        chosen_layout = get_layout(slide_spec.get("layoutId"))
        new_slide = prs.slides.add_slide(chosen_layout)
        fill_slide(new_slide, slide_spec)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()
